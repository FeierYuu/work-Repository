#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
演示文稿生成模块
负责将分析后的内容转换为PPTX格式的演示文稿
"""

from typing import Dict, List, Any
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from content_analyzer import ContentAnalyzer
from style_optimizer import StyleOptimizer


class PresentationGenerator:
    """
    演示文稿生成器类
    """
    
    def __init__(self):
        self.analyzer = ContentAnalyzer()
        self.style_optimizer = StyleOptimizer()
        # 预定义数学公式的字体设置
        self.formula_font_name = 'Courier New'
        self.formula_font_size = Pt(32)
    
    def generate(self, content_blocks: List[Dict[str, Any]], output_path: str, 
                style_options: Dict[str, Any] = None) -> str:
        """
        生成演示文稿
        
        Args:
            content_blocks: 分析后的内容块列表
            output_path: 输出文件路径
            style_options: 样式优化选项
            
        Returns:
            str: 生成的文件路径
        """
        print(f"\n========== 开始生成演示文稿 ==========")
        print(f"输出路径: {output_path}")
        print(f"内容块总数: {len(content_blocks)}")
        
        # 输出内容块详细信息
        print("\n内容块详细信息:")
        for i, block in enumerate(content_blocks):
            block_type = block.get('type', 'unknown')
            if block_type == 'image':
                path = block.get('path', 'no_path')
                print(f"  块{i+1}: 类型={block_type}, 路径={path}")
            else:
                print(f"  块{i+1}: 类型={block_type}")
        
        # 创建演示文稿对象
        print("\n创建新的演示文稿对象...")
        prs = Presentation()
        print(f"演示文稿对象创建成功")
        print(f"幻灯片尺寸: {prs.slide_width} x {prs.slide_height}")
        
        # 为每个内容块生成幻灯片
        print("\n开始处理内容块...")
        for i, block in enumerate(content_blocks):
            block_type = block.get('type', 'unknown')
            print(f"\n处理内容块 {i+1}/{len(content_blocks)}: 类型={block_type}")
            
            if block_type == 'section':
                print("生成章节幻灯片...")
                self._generate_section_slides(prs, block)
                print(f"章节幻灯片生成完成")
            elif block_type == 'image':
                print(f"生成图片幻灯片: {block.get('path', '无路径')}")
                self._generate_slide_with_image(prs, block)
                print(f"图片幻灯片生成完成")
            else:
                print(f"生成内容幻灯片...")
                self._generate_content_slide(prs, block)
                print(f"内容幻灯片生成完成")
        
        # 检查生成的幻灯片数量
        print(f"\n所有内容块处理完成，共生成 {len(prs.slides)} 张幻灯片")
        
        # 应用样式优化
        print(f"应用样式优化...")
        self.style_optimizer.optimize_presentation(prs, style_options)
        print(f"样式优化完成")
        
        # 保存演示文稿
        print(f"\n保存演示文稿到: {output_path}")
        try:
            # 确保输出目录存在
            import os
            output_dir = os.path.dirname(output_path)
            if output_dir and not os.path.exists(output_dir):
                print(f"创建输出目录: {output_dir}")
                os.makedirs(output_dir)
            
            prs.save(output_path)
            print(f"✓ 演示文稿保存成功")
            
            # 验证文件是否存在
            if os.path.exists(output_path):
                file_size = os.path.getsize(output_path)
                print(f"生成的文件大小: {file_size} 字节")
            else:
                print(f"✗ 警告: 无法验证文件是否保存成功")
                
        except Exception as save_err:
            print(f"✗ 保存演示文稿失败: {str(save_err)}")
            import traceback
            traceback.print_exc()
            raise
        
        print(f"\n========== 演示文稿生成结束 ==========")
        return output_path
    
    def _generate_section_slides(self, prs: Presentation, section: Dict[str, Any]):
        """
        为一个章节生成幻灯片
        
        Args:
            prs: 演示文稿对象
            section: 章节内容块
        """
        # 创建章节标题幻灯片
        title_slide_layout = prs.slide_layouts[0]  # 标题幻灯片
        slide = prs.slides.add_slide(title_slide_layout)
        
        # 设置标题
        title = slide.shapes.title
        title.text = section['title']
        title.text_frame.paragraphs[0].font.size = Pt(44)  # 大字体
        title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # 可以添加副标题或其他信息
        subtitle = slide.placeholders[1]
        subtitle.text = ""
        
        # 为章节内容生成幻灯片
        for content_item in section['content']:
            self._generate_content_slide(prs, content_item, section['title'])
    
    def _generate_content_slide(self, prs: Presentation, content: Dict[str, Any], 
                               section_title: str = None):
        """
        为内容项生成幻灯片
        
        Args:
            prs: 演示文稿对象
            content: 内容项
            section_title: 所属章节标题
        """
        if content['type'] == 'paragraph':
            # 为段落内容创建幻灯片
            content_slide_layout = prs.slide_layouts[1]  # 标题和内容
            slide = prs.slides.add_slide(content_slide_layout)
            
            # 设置标题
            title = slide.shapes.title
            if content.get('title'):
                title.text = content['title']
            elif section_title:
                title.text = section_title
            else:
                title.text = "内容"
            
            title.text_frame.paragraphs[0].font.size = Pt(32)  # 大字体
            title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # 设置内容
            content_shape = slide.placeholders[1]
            tf = content_shape.text_frame
            
            # 清空默认内容
            tf.clear()
            
            # 添加摘要文本
            summarized_text = self.analyzer.summarize_text(content['content'], max_length=500)
            p = tf.add_paragraph()
            p.text = summarized_text
            p.font.size = Pt(24)  # 大字体
            p.alignment = PP_ALIGN.LEFT
            
        elif content['type'] == 'table':
            # 为表格创建幻灯片
            content_slide_layout = prs.slide_layouts[1]  # 标题和内容
            slide = prs.slides.add_slide(content_slide_layout)
            
            # 设置标题
            title = slide.shapes.title
            title.text = content.get('title', '表格')
            title.text_frame.paragraphs[0].font.size = Pt(32)  # 大字体
            title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # 创建表格
            table_data = content['content']
            if table_data:
                rows, cols = len(table_data), len(table_data[0])
                
                # 设置表格位置和大小
                left = Inches(1)
                top = Inches(2)
                width = Inches(11)
                height = Inches(4)
                
                # 添加表格
                table = slide.shapes.add_table(rows, cols, left, top, width, height).table
                
                # 填充表格数据
                for r, row_data in enumerate(table_data):
                    for c, cell_data in enumerate(row_data):
                        cell = table.cell(r, c)
                        cell.text = cell_data
                        
                        # 设置单元格文本格式
                        cell.text_frame.paragraphs[0].font.size = Pt(18)  # 大字体
                        
                        # 表头样式
                        if r == 0:
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = RGBColor(200, 200, 200)
                            cell.text_frame.paragraphs[0].font.bold = True
        
        elif content['type'] == 'formula':
            # 为公式创建幻灯片
            content_slide_layout = prs.slide_layouts[1]  # 标题和内容
            slide = prs.slides.add_slide(content_slide_layout)
            
            # 设置标题
            title = slide.shapes.title
            title.text = content.get('title', '公式')
            title.text_frame.paragraphs[0].font.size = Pt(32)  # 大字体
            title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # 添加公式内容
            content_shape = slide.placeholders[1]
            tf = content_shape.text_frame
            tf.clear()
            
            # 检查是否为LaTeX公式
            is_latex = content.get('is_latex', False)
            
            # 为公式创建段落
            p = tf.add_paragraph()
            
            # 对于LaTeX公式，去除$符号，使其更易读
            formula_content = content['content']
            if is_latex:
                # 去除$符号，保留公式本身
                formula_content = formula_content.replace('$', '')
                # 添加LaTeX标记到标题
                title.text = title.text + " (LaTeX格式)"
            
            p.text = formula_content
            # 使用专门的公式字体设置
            p.font.name = self.formula_font_name
            p.font.size = self.formula_font_size
            p.alignment = PP_ALIGN.CENTER
            
            # 应用样式优化器来增强公式显示
            self.style_optimizer.optimize_formula_display(tf)
    
    def _generate_slide_with_image(self, prs: Presentation, block: Dict[str, Any]):
        """
        生成包含图片的幻灯片
        
        Args:
            prs: 演示文稿对象
            block: 图片内容块
        """
        print(f"\n========== 开始生成图片幻灯片 ==========")
        import os
        import io
        from PIL import Image
        
        # 从block中提取所有关键信息
        image_path = block.get('path', '')
        image_title = block.get('title', '')
        image_caption = block.get('caption', '')
        print(f"图片路径: {image_path}")
        print(f"图片标题: {image_title}")
        print(f"图片说明: {image_caption}")
        
        # 验证图片文件是否存在
        if not image_path:
            print(f"错误: 图片路径为空")
        elif not os.path.exists(image_path):
            print(f"错误: 图片文件不存在: {image_path}")
        else:
            # 验证图片文件是否有效
            try:
                # 检查文件大小
                file_size = os.path.getsize(image_path)
                print(f"图片文件大小: {file_size} 字节")
                
                if file_size == 0:
                    print(f"警告: 图片文件为空: {image_path}")
            except Exception as file_err:
                print(f"检查图片文件失败: {str(file_err)}")
        
        try:
            # 使用带有标题和内容的幻灯片布局
            print("获取幻灯片布局...")
            content_slide_layout = prs.slide_layouts[1]  # 标题和内容
            print("添加新幻灯片...")
            slide = prs.slides.add_slide(content_slide_layout)
            print(f"幻灯片添加成功，当前幻灯片数量: {len(prs.slides)}")
            
            # 设置标题
            print("设置幻灯片标题...")
            title = slide.shapes.title
            
            # 提取图片文件名作为标题
            if image_path:
                image_name = os.path.basename(image_path)
                title.text = image_title if image_title else f"图片: {image_name}"
                print(f"标题设置为: {title.text}")
            else:
                title.text = image_title if image_title else "图片"
                print(f"标题设置为: {title.text}")
            
            title.text_frame.paragraphs[0].font.size = Pt(32)
            title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            print("标题格式设置完成")
            
            # 获取内容占位符
            print("查找内容占位符...")
            content_placeholder = None
            for shape in slide.placeholders:
                print(f"  占位符ID: {shape.placeholder_format.idx}, 类型: {shape.placeholder_format.type}")
                if shape.placeholder_format.idx == 1:  # 内容占位符的索引通常是1
                    content_placeholder = shape
                    print(f"  找到内容占位符: {content_placeholder}")
                    break
            
            # 确保内容占位符被清空，避免与图片重叠
            if content_placeholder:
                print("清空内容占位符...")
                content_placeholder.text = ""
            else:
                print("警告: 未找到内容占位符")
            
            # 添加图片 - 修复图片显示问题
            if image_path and os.path.exists(image_path):
                print(f"\n开始处理图片: {image_path}")
                print(f"文件大小: {os.path.getsize(image_path)} 字节")
                
                # 计算图片大小和位置
                slide_width = prs.slide_width
                slide_height = prs.slide_height
                print(f"幻灯片尺寸: {slide_width} x {slide_height}")
                
                # 使用更好的边距计算
                left_margin = Inches(1.0)
                top_margin = Inches(2.5)  # 给标题留出足够空间
                print(f"边距设置: 左={left_margin}, 上={top_margin}")
                
                # 计算可用空间
                available_width = slide_width - left_margin * 2
                available_height = slide_height - top_margin - Inches(1.0)
                print(f"可用空间: {available_width} x {available_height}")
                
                new_width = None
                new_height = None
                left = None
                top = None
                
                try:
                    # 打开图片文件
                    print("打开图片文件进行分析...")
                    with Image.open(image_path) as img:
                        width, height = img.size
                        img_format = img.format
                        img_mode = img.mode
                        print(f"图片信息: 尺寸={width}x{height} 像素, 格式={img_format}, 模式={img_mode}")
                        
                        # 计算缩放比例
                        width_ratio = available_width / width
                        height_ratio = available_height / height
                        scale = min(width_ratio, height_ratio, 1.0)  # 不超过原始大小
                        print(f"缩放计算: 宽度比例={width_ratio:.2f}, 高度比例={height_ratio:.2f}, 最终比例={scale:.2f}")
                        
                        # 计算新尺寸（使用Inches单位确保兼容性）
                        new_width = Inches(width * scale / 96)  # 假设DPI为96
                        new_height = Inches(height * scale / 96)
                        print(f"缩放后尺寸: {new_width}x{new_height}")
                        
                        # 居中图片
                        left = (slide_width - new_width) / 2
                        top = top_margin
                        print(f"计算图片位置: 左={left}, 上={top}")
                except Exception as img_err:
                    print(f"图片处理错误: {img_err}")
                    import traceback
                    traceback.print_exc()
                    # 如果图片读取失败，使用默认尺寸
                    print("使用默认尺寸作为备选...")
                    new_width = Inches(6.0)
                    new_height = Inches(4.0)
                    left = (slide_width - new_width) / 2
                    top = top_margin
                    print(f"默认尺寸: {new_width}x{new_height}, 默认位置: 左={left}, 上={top}")
                
                # 使用BytesIO缓冲区读取图片数据，避免文件路径问题
                print("\n准备添加图片到幻灯片...")
                try:
                    # 读取图片文件到内存
                    print("读取图片文件到内存缓冲区...")
                    with open(image_path, 'rb') as f:
                        image_data = f.read()
                    print(f"图片数据读取成功，数据大小: {len(image_data)} 字节")
                    
                    # 创建BytesIO缓冲区
                    image_stream = io.BytesIO(image_data)
                    print(f"创建内存缓冲区成功")
                    
                    # 尝试多种添加图片的方法
                    added_successfully = False
                    picture = None
                    
                    # 方法1: 使用内存数据和计算的尺寸
                    print("\n尝试方法1: 使用内存数据和计算的尺寸")
                    try:
                        picture = slide.shapes.add_picture(
                            image_stream,
                            left=left,
                            top=top,
                            width=new_width,
                            height=new_height
                        )
                        print(f"✓ 方法1成功: 图片添加到幻灯片，对象ID: {id(picture)}")
                        print(f"图片位置: left={left}, top={top}")
                        added_successfully = True
                    except Exception as method1_err:
                        print(f"✗ 方法1失败: {method1_err}")
                    
                    # 如果方法1失败，尝试方法2: 不指定尺寸
                    if not added_successfully:
                        print("\n尝试方法2: 使用内存数据但不指定尺寸")
                        try:
                            image_stream.seek(0)  # 重置文件指针
                            picture = slide.shapes.add_picture(
                                image_stream,
                                left=left,
                                top=top
                            )
                            print(f"✓ 方法2成功: 图片添加到幻灯片，对象ID: {id(picture)}")
                            added_successfully = True
                        except Exception as method2_err:
                            print(f"✗ 方法2失败: {method2_err}")
                    
                    # 如果方法2失败，尝试方法3: 直接使用文件路径
                    if not added_successfully:
                        print("\n尝试方法3: 直接使用文件路径")
                        try:
                            picture = slide.shapes.add_picture(
                                image_path,
                                left=left,
                                top=top,
                                width=new_width,
                                height=new_height
                            )
                            print(f"✓ 方法3成功: 图片添加到幻灯片，对象ID: {id(picture)}")
                            added_successfully = True
                        except Exception as method3_err:
                            print(f"✗ 方法3失败: {method3_err}")
                    
                    # 如果方法3失败，尝试方法4: 使用文件路径但不指定尺寸
                    if not added_successfully:
                        print("\n尝试方法4: 使用文件路径但不指定尺寸")
                        try:
                            picture = slide.shapes.add_picture(
                                image_path,
                                left=left,
                                top=top
                            )
                            print(f"✓ 方法4成功: 图片添加到幻灯片，对象ID: {id(picture)}")
                            added_successfully = True
                        except Exception as method4_err:
                            print(f"✗ 方法4失败: {method4_err}")
                    
                    # 检查是否成功添加图片
                    if added_successfully and picture:
                        print(f"\n图片添加成功! 当前幻灯片中的形状数量: {len(slide.shapes)}")
                        
                        # 添加图片说明文本
                        if image_caption:
                            print(f"添加图片说明: {image_caption}")
                            caption_shape = slide.shapes.add_textbox(
                                left=left,
                                top=top + new_height + Inches(0.2),
                                width=new_width,
                                height=Inches(0.5)
                            )
                            tf = caption_shape.text_frame
                            tf.text = image_caption
                            tf.paragraphs[0].font.size = Pt(14)
                            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
                            print("图片说明添加成功")
                    else:
                        print(f"✗ 所有添加图片的方法都失败了")
                        # 添加错误信息
                        if content_placeholder:
                            content_placeholder.text = "图片添加失败: 所有尝试的方法都未成功"
                except Exception as add_err:
                    print(f"添加图片过程中出错: {add_err}")
                    import traceback
                    traceback.print_exc()
                    # 添加错误信息
                    if content_placeholder:
                        content_placeholder.text = f"图片添加失败: {str(add_err)}"
            else:
                print(f"警告: 无法找到图片文件: {image_path}")
                if content_placeholder:
                    content_placeholder.text = "图片无法加载: 文件不存在或路径无效"
        
        except Exception as e:
            print(f"添加图片到幻灯片时出错: {str(e)}")
            import traceback
            traceback.print_exc()
            # 如果出错，添加一个错误说明文本
            for shape in slide.placeholders:
                if shape.placeholder_format.idx == 1:
                    shape.text = f"添加图片时出错: {str(e)}"
                    break
        finally:
            print(f"========== 图片幻灯片生成结束 ==========\n")
    
    def optimize_slides(self, prs: Presentation, options: Dict[str, Any] = None):
        """
        优化幻灯片布局和样式
        
        Args:
            prs: 演示文稿对象
            options: 样式优化选项
        """
        # 使用样式优化器进行优化
        self.style_optimizer.optimize_presentation(prs, options)


if __name__ == "__main__":
    # 测试代码
    generator = PresentationGenerator()
    print("演示文稿生成模块已创建")