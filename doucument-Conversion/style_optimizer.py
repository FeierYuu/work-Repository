#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
排版和样式优化模块
负责优化演示文稿的字体、颜色、布局等样式
"""

from typing import Dict, List, Any, Tuple
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE


class StyleOptimizer:
    """
    样式优化器类
    负责优化演示文稿的排版和样式
    """
    
    def __init__(self):
        # 定义颜色主题
        self.theme_colors = {
            'title': RGBColor(31, 73, 125),  # 深蓝色标题
            'background': RGBColor(255, 255, 255),  # 白色背景
            'text': RGBColor(0, 0, 0),  # 黑色文本
            'accent': RGBColor(227, 114, 34),  # 橙色强调
            'table_header': RGBColor(51, 92, 143),  # 表格头部背景
            'table_header_text': RGBColor(255, 255, 255),  # 表格头部文本
        }
        
        # 定义字体大小
        self.font_sizes = {
            'slide_title': Pt(44),  # 幻灯片标题
            'section_title': Pt(36),  # 节标题
            'content_title': Pt(32),  # 内容标题
            'normal_text': Pt(24),  # 普通文本
            'small_text': Pt(20),  # 小文本
            'table_text': Pt(18),  # 表格文本
        }
    
    def optimize_presentation(self, prs: Presentation, options: Dict[str, Any] = None):
        """
        优化整个演示文稿的样式
        
        Args:
            prs: 演示文稿对象
            options: 优化选项
        """
        if options is None:
            options = {}
        
        # 优化幻灯片大小（确保横向）
        self._optimize_slide_size(prs)
        
        # 为每个幻灯片应用样式
        for slide in prs.slides:
            self._optimize_slide(slide, options)
        
        return prs
    
    def _optimize_slide_size(self, prs: Presentation):
        """
        优化幻灯片大小（确保横向）
        
        Args:
            prs: 演示文稿对象
        """
        # 设置为标准宽屏横向尺寸 (16:9)
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
    
    def _optimize_slide(self, slide, options: Dict[str, Any]):
        """
        优化单个幻灯片的样式
        
        Args:
            slide: 幻灯片对象
            options: 优化选项
        """
        # 优化标题样式
        if slide.shapes.title:
            self._optimize_title(slide.shapes.title)
        
        # 优化内容形状
        for shape in slide.shapes:
            if shape.has_text_frame:
                self._optimize_text_frame(shape.text_frame, options)
            elif shape.has_table:
                self._optimize_table(shape.table)
    
    def _optimize_title(self, title_shape):
        """
        优化标题样式
        
        Args:
            title_shape: 标题形状对象
        """
        text_frame = title_shape.text_frame
        
        # 设置标题文本格式
        for paragraph in text_frame.paragraphs:
            paragraph.font.size = self.font_sizes['slide_title']
            paragraph.font.color.rgb = self.theme_colors['title']
            paragraph.font.name = '微软雅黑'
            paragraph.font.bold = True
            paragraph.alignment = PP_ALIGN.CENTER
            
            # 设置行间距
            paragraph.space_after = Pt(12)
    
    def _optimize_text_frame(self, text_frame, options: Dict[str, Any]):
        """
        优化文本框样式
        
        Args:
            text_frame: 文本框对象
            options: 优化选项
        """
        # 自动调整文本框大小
        text_frame.word_wrap = True
        
        # 优化每个段落
        for i, paragraph in enumerate(text_frame.paragraphs):
            # 根据段落位置和内容长度调整字体大小
            if i == 0 and len(paragraph.text) < 30:  # 可能是副标题
                paragraph.font.size = self.font_sizes['content_title']
                paragraph.font.color.rgb = self.theme_colors['title']
                paragraph.font.bold = True
            else:
                # 根据内容长度调整字体大小
                if len(paragraph.text) > 100:
                    paragraph.font.size = self.font_sizes['small_text']
                else:
                    paragraph.font.size = self.font_sizes['normal_text']
                
                paragraph.font.color.rgb = self.theme_colors['text']
                
            # 设置字体
            paragraph.font.name = '微软雅黑'
            
            # 设置对齐方式
            if len(paragraph.text) > 50:
                paragraph.alignment = PP_ALIGN.LEFT
            else:
                paragraph.alignment = PP_ALIGN.CENTER
            
            # 设置行间距
            paragraph.space_after = Pt(6)
    
    def _optimize_table(self, table):
        """
        优化表格样式
        
        Args:
            table: 表格对象
        """
        # 优化表头
        if table.rows:
            header_row = table.rows[0]
            for cell in header_row.cells:
                # 设置表头背景色
                cell.fill.solid()
                cell.fill.fore_color.rgb = self.theme_colors['table_header']
                
                # 设置表头文本格式
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.color.rgb = self.theme_colors['table_header_text']
                    paragraph.font.size = self.font_sizes['table_text']
                    paragraph.font.bold = True
                    paragraph.alignment = PP_ALIGN.CENTER
        
        # 优化表格内容
        if len(table.rows) > 1:
            for row_idx in range(1, len(table.rows)):
                row = table.rows[row_idx]
                for cell in row.cells:
                    # 斑马纹效果
                    if (row_idx - 1) % 2 == 1:  # 从第二行开始，索引从0重新计算
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(240, 240, 240)
                    else:
                        cell.fill.background()
                    
                    # 设置文本格式
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.size = self.font_sizes['table_text']
                        paragraph.font.name = '微软雅黑'
                        paragraph.alignment = PP_ALIGN.CENTER
    
    def ensure_text_visibility(self, text_frame):
        """
        确保文本可见性，调整文本大小以适应文本框
        
        Args:
            text_frame: 文本框对象
        """
        # 简单的文本适配逻辑
        # 实际应用可能需要更复杂的计算
        total_text_length = sum(len(paragraph.text) for paragraph in text_frame.paragraphs)
        
        # 根据文本总量调整字体大小
        if total_text_length > 500:
            target_size = Pt(20)
        elif total_text_length > 300:
            target_size = Pt(22)
        else:
            target_size = Pt(24)
        
        # 应用调整后的字体大小
        for paragraph in text_frame.paragraphs:
            if paragraph.font.size > target_size:
                paragraph.font.size = target_size
    
    def add_visual_elements(self, slide, options: Dict[str, Any] = None):
        """
        添加视觉元素增强幻灯片
        
        Args:
            slide: 幻灯片对象
            options: 配置选项
        """
        if options is None:
            options = {}
            
        # 可以添加一些简单的装饰元素
        if options.get('add_decorations', False):
            # 添加一个简单的装饰条
            left = Inches(0)
            top = Inches(7)
            width = Inches(13.33)
            height = Inches(0.5)
            
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                left, top, width, height
            )
            
            # 设置装饰条颜色
            fill = shape.fill
            fill.solid()
            fill.fore_color.rgb = self.theme_colors['title']
            
            # 隐藏边框
            line = shape.line
            line.fill.background()
    
    def optimize_formula_display(self, text_frame):
        """
        优化公式显示
        
        Args:
            text_frame: 包含公式的文本框
        """
        # 为公式文本设置等宽字体，提高可读性
        for paragraph in text_frame.paragraphs:
            # 检测是否包含公式符号
            has_formula_chars = any(char in paragraph.text for char in ['\\', '=', '+', '-', '*', '/', '∑', '∫', 'π', '√', '^', '_', '{', '}'])
            
            if has_formula_chars:
                # 设置等宽字体以确保公式对齐正确
                paragraph.font.name = 'Courier New'  # 等宽字体
                paragraph.font.size = Pt(32)  # 更大的字体
                paragraph.alignment = PP_ALIGN.CENTER
                
                # 设置公式文本为粗体，增强视觉效果
                paragraph.font.bold = True
                
                # 为公式文本设置专业的颜色
                paragraph.font.color.rgb = RGBColor(0, 51, 102)  # 深蓝色，更专业
                
                # 增加行距，使公式更易读
                paragraph.space_after = Pt(6)
                
                # 调整段落缩进
                paragraph.left_indent = Pt(0)
                paragraph.right_indent = Pt(0)
                paragraph.space_before = Pt(12)  # 增加上方空间


# 更新演示文稿生成器以使用样式优化器
# 这个函数用于将样式优化器集成到演示文稿生成流程中
def update_presentation_generator():
    """
    更新演示文稿生成器代码，集成样式优化器
    """
    # 这个函数只是作为一个提示，实际集成需要修改 presentation_generator.py 文件
    pass


if __name__ == "__main__":
    # 测试代码
    optimizer = StyleOptimizer()
    print("样式优化模块已创建")