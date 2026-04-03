import os
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Inches
from PIL import Image

# ===================== 配置项 =====================
HTML_FILES = [
    "Системы ИИ AI Systems.html",
    "ch1.html",
    "ch2.html",
    "Logic_programming.html"
]
PAPKA_IMG = "slides_final"  # 截图保存目录
SHIRINA = 1920             # 窗口宽度（防遮挡）
VYSOTA = 1080              # 窗口高度（防遮挡）
# ==================================================

def snimok_slaydov(put_html):
    # 截取单个HTML演示文稿所有页面
    imya_bez_rasshir = os.path.splitext(put_html)[0]
    papka_slaydov = os.path.join(PAPKA_IMG, imya_bez_rasshir)
    os.makedirs(papka_slaydov, exist_ok=True)

    with sync_playwright() as p:
        brauzer = p.chromium.launch(headless=True, args=["--no-sandbox"])
        stranitsa = brauzer.new_page(viewport={"width": SHIRINA, "height": VYSOTA})
        stranitsa.goto(f"file://{os.path.abspath(put_html)}")
        stranitsa.wait_for_timeout(3000)

        # 获取总页数
        try:
            vsego_slaydov = stranitsa.evaluate("document.querySelectorAll('.slide').length")
        except:
            vsego_slaydov = 20

        print(f"📄 Обработка: {put_html} | Всего слайдов: {vsego_slaydov}")
        izobrazheniya = []

        for nomer in range(vsego_slaydov):
            put_kartinki = os.path.join(papka_slaydov, f"slayd_{nomer+1:02d}.png")
            stranitsa.screenshot(path=put_kartinki)
            izobrazheniya.append(put_kartinki)
            print(f"   Слайд {nomer+1} готов")

            # 翻下一页
            if nomer < vsego_slaydov - 1:
                try:
                    stranitsa.evaluate("changeSlide(1)")
                except:
                    pass
                stranitsa.wait_for_timeout(500)

        brauzer.close()
    return izobrazheniya, imya_bez_rasshir

def eksport_v_pdf_pptx(izobrazheniya, imya):
    # 生成PPTX文件
    prezentaciya = Presentation()
    prezentaciya.slide_width = Inches(13.33)
    prezentaciya.slide_height = Inches(7.5)
    pustoy_maket = prezentaciya.slide_layouts[6]

    for kart in izobrazheniya:
        slayd = prezentaciya.slides.add_slide(pustoy_maket)
        slayd.shapes.add_picture(kart, left=Inches(0), top=Inches(0),
                                 width=prezentaciya.slide_width, height=prezentaciya.slide_height)
    prezentaciya.save(f"{imya}.pptx")

    # 生成PDF文件
    karty = [Image.open(kart).convert("RGB") for kart in izobrazheniya]
    karty[0].save(f"{imya}.pdf", save_all=True, append_images=karty[1:])
    print(f"✅ Экспорт завершен: {imya}.pdf + {imya}.pptx\n")

# ===================== 主程序 =====================
if __name__ == "__main__":
    print("🚀 Запуск пакетной конвертации презентаций\n")
    os.makedirs(PAPKA_IMG, exist_ok=True)

    for html in HTML_FILES:
        if os.path.exists(html):
            imgs, name = snimok_slaydov(html)
            eksport_v_pdf_pptx(imgs, name)
        else:
            print(f"❌ Файл не найден: {html}")

    print("🎉 Все операции выполнены успешно!")